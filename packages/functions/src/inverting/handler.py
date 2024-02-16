from typing import List
import os
import json
from pptx import Presentation
from io import BytesIO
import boto3
from pydantic import BaseModel
from aws_lambda_powertools.utilities.parser import (
    ValidationError,
    event_parser,
    parse,
)
from aws_lambda_powertools.utilities.parser.models import APIGatewayProxyEventV2Model
from aws_lambda_powertools.utilities.typing import LambdaContext
from pptx.dml.color import RGBColor
from zipfile import ZipFile, ZIP_DEFLATED
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image, ImageOps
from aws_lambda_powertools import Logger

BUCKET = os.environ["BUCKET"]

logger = Logger()
s3 = boto3.client("s3")


class InvertRequest(BaseModel):
    file_keys: List[str]


# @event_parser(model=InvertRequest)
def handler(event, context: LambdaContext):
    try:
        # Check if event is bytes
        if isinstance(event, bytes):
            # Decode bytes string and convert to dictionary
            event = json.loads(event.decode("utf-8"))

        # Parse dictionary to Pydantic model
        invert_request = InvertRequest(**event)
        zip_buffer = BytesIO()
        with ZipFile(zip_buffer, "a", ZIP_DEFLATED, False) as zip_file:
            for file_key in invert_request.file_keys:
                # event = parse(event=json.loads(event.decode("utf-8")), model=InvertRequest)
                response = s3.get_object(Bucket=BUCKET, Key=file_key)
                body = response["Body"].read()

                presentation = Presentation(BytesIO(body))

                for slide in presentation.slides:
                    # Invert slide background to black
                    background = slide.background
                    fill = background.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor(0, 0, 0)

                    # Invert text color to white for all shapes that contain text
                    for shape in slide.shapes:
                        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            # Extract the image from the shape
                            image_stream = BytesIO(shape.image.blob)
                            with Image.open(image_stream) as img:
                                # If image format is not supported convert it first
                                if img.mode not in ("RGB", "RGBA"):
                                    img = img.convert("RGB")

                                # Invert the image color using Pillow only if supported
                                inverted_img = ImageOps.invert(img)

                                # Save the inverted image to a stream
                                img_stream = BytesIO()
                                inverted_img.save(img_stream, format="PNG")
                                img_stream.seek(0)

                                # Remove the original picture
                                left = shape.left
                                top = shape.top
                                width = shape.width
                                height = shape.height
                                slide.shapes._spTree.remove(shape._element)

                                # Add the new inverted image to the slide
                                slide.shapes.add_picture(
                                    img_stream, left, top, width, height
                                )

                        if not shape.has_text_frame:
                            continue

                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.color.rgb = RGBColor(255, 255, 255)

                output = BytesIO()
                presentation.save(output)
                output.seek(0)
                new_key = "inverted/" + file_key
                zip_file.writestr(new_key, output.getvalue())

        zip_buffer.seek(0)
        key = "inverted.zip"
        s3.put_object(Bucket=BUCKET, Key=key, Body=zip_buffer.getvalue())
        print(f"Uploaded inverted folder to {key}")
        return {"status": "success", "key": key}
    except ValidationError as e:
        print(e)
        return {"error": e.errors()}

    except Exception as e:
        print(e)
        return {"error": str(e)}
