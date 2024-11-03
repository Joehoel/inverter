#!/usr/bin/env python3

import os
import sys
import signal
import argparse
import time
import hashlib
from pathlib import Path
from typing import List, Tuple, Dict
from io import BytesIO
from concurrent.futures import ThreadPoolExecutor, as_completed, TimeoutError
import threading
import contextlib
from PIL import Image, ImageOps
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from rich.progress import Progress, BarColumn, TextColumn, TimeRemainingColumn
from rich.console import Console

from rich.traceback import install

install(show_locals=True)

# Global flags for interrupt handling
interrupted = False
force_exit = False
interrupt_count = 0
interrupt_lock = threading.Lock()
last_interrupt_time = 0

console = Console()


def signal_handler(signum, frame):
    """Handle keyboard interrupt gracefully with force exit option"""
    global interrupted, force_exit, interrupt_count, last_interrupt_time
    current_time = time.time()

    with interrupt_lock:
        if current_time - last_interrupt_time < 1:  # Double interrupt within 1 second
            console.print("\nForce exiting...", style="bold red")
            os._exit(1)  # Force exit the program

        interrupt_count += 1
        last_interrupt_time = current_time

        if not interrupted:
            console.print(
                "\nInterrupt received. Finishing current files and cleaning up...",
                style="bold yellow",
            )
            console.print(
                "Press Ctrl+C again within 1 second to force exit.", style="bold yellow"
            )
            interrupted = True


# Set up signal handler
signal.signal(signal.SIGINT, signal_handler)


def calculate_file_hash(file_path: Path) -> str:
    """Calculate SHA-256 hash of a file"""
    sha256_hash = hashlib.sha256()
    sha256_hash.update(file_path.name.encode("utf-8"))
    return sha256_hash.hexdigest()


def get_existing_files_info(
    directory: Path, recursive: bool = False
) -> Dict[Path, str]:
    """Get information about existing files in the output directory"""
    pattern = "**/*.pptx" if recursive else "*.pptx"
    existing_files = {}

    if directory.exists():
        for file_path in directory.glob(pattern):
            try:
                existing_files[file_path] = calculate_file_hash(file_path)
            except Exception:
                continue

    return existing_files


def should_process_file(
    input_path: Path, output_path: Path, existing_files: Dict[Path, str]
) -> bool:
    """Determine if a file needs to be processed"""
    if not output_path.exists():
        return True

    try:
        input_hash = calculate_file_hash(input_path)
        output_hash = existing_files.get(output_path)
        return not output_hash or input_hash != output_hash
    except Exception:
        return True


@contextlib.contextmanager
def managed_image(image_stream):
    """Context manager to ensure proper cleanup of image resources"""
    img = Image.open(image_stream)
    try:
        yield img
    finally:
        img.close()


def process_image(shape, slide):
    """Process and invert a single image shape"""
    image_stream = BytesIO(shape.image.blob)
    with managed_image(image_stream) as img:
        if img.mode not in ("RGB", "RGBA"):
            img = img.convert("RGB")
        inverted_img = ImageOps.invert(img)
        img_stream = BytesIO()
        inverted_img.save(img_stream, format="PNG")
        img_stream.seek(0)

        # Store shape properties before removal
        left, top = shape.left, shape.top
        width, height = shape.width, shape.height

        # Remove old shape and add new one
        slide.shapes._spTree.remove(shape._element)
        slide.shapes.add_picture(img_stream, left, top, width, height)


def invert_presentation(bytes_data):
    """Inverts colors in a PowerPoint presentation"""
    presentation = Presentation(bytes_data)

    for slide in presentation.slides:
        # Invert slide background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)

        # Process shapes
        shapes_to_process = list(slide.shapes)  # Create a copy of the shapes list
        for shape in shapes_to_process:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                process_image(shape, slide)
            elif shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = RGBColor(255, 255, 255)

    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def get_default_output_dir(input_path: Path) -> Path:
    """Generate default output directory path with (inverted) suffix"""
    return input_path.parent / f"{input_path.name} (inverted)"


def process_file_with_timeout(
    input_path: Path, output_path: Path, timeout: int = 300
) -> Tuple[Path, bool, str]:
    """
    Process a single PowerPoint file with timeout.
    Returns tuple of (input_path, success_boolean, error_message).
    """
    try:
        # Check for interrupt
        global interrupted
        if interrupted:
            return input_path, False, "Interrupted by user"

        # Ensure output directory exists
        output_path.parent.mkdir(parents=True, exist_ok=True)

        # Read and process the file
        with open(input_path, "rb") as f:
            inverted_pptx = invert_presentation(f)

        # Save the inverted presentation
        with open(output_path, "wb") as f:
            f.write(inverted_pptx)

        return input_path, True, ""

    except Exception as e:
        return input_path, False, str(e)


def find_pptx_files(directory: Path, recursive: bool = False) -> List[Path]:
    """Find all .pptx files in directory"""
    pattern = "**/*.pptx" if recursive else "*.pptx"
    return list(directory.glob(pattern))


def process_directory(
    input_dir: Path,
    output_dir: Path,
    recursive: bool = False,
    max_workers: int = None,
    timeout: int = 300,
) -> tuple[int, int, int]:
    """
    Process all PowerPoint files in a directory using parallel processing.
    Returns tuple of (success_count, failure_count, skipped_count).
    """
    global interrupted
    success_count = 0
    failure_count = 0
    skipped_count = 0

    # Create output directory if it doesn't exist
    output_dir.mkdir(parents=True, exist_ok=True)

    # Get information about existing files
    existing_files = get_existing_files_info(output_dir, recursive)

    # Collect all .pptx files
    pptx_files = find_pptx_files(input_dir, recursive)

    if not pptx_files:
        console.print(f"No PowerPoint files found in {input_dir}", style="bold red")
        return 0, 0, 0

    # Prepare the work items - maintaining directory structure
    work_items = []
    for input_path in pptx_files:
        # Calculate relative path from input directory
        rel_path = input_path.relative_to(input_dir)
        # Use same relative path in output directory
        output_path = output_dir / rel_path

        # Check if we need to process this file
        if should_process_file(input_path, output_path, existing_files):
            work_items.append((input_path, output_path))
        else:
            skipped_count += 1

    # If no files need processing, return early
    if not work_items:
        console.print("All files are up to date.", style="bold green")
        return 0, 0, skipped_count

    # Calculate optimal number of workers if not specified
    if max_workers is None:
        max_workers = min(32, (os.cpu_count() or 1) * 4)

    # Create progress bar
    total_files = len(work_items)

    with Progress(
        TextColumn("[progress.description]{task.description}"),
        BarColumn(),
        "[progress.percentage]{task.percentage:>3.0f}%",
        "•",
        TimeRemainingColumn(),
        transient=True,
    ) as progress:
        task = progress.add_task("Processing presentations", total=total_files)

        # Process files in parallel
        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            # Submit all tasks
            future_to_path = {
                executor.submit(
                    process_file_with_timeout, input_path, output_path
                ): input_path
                for input_path, output_path in work_items
            }

            # Process completed tasks
            for future in as_completed(future_to_path):
                if interrupted:
                    executor.shutdown(wait=False, cancel_futures=True)
                    break

                try:
                    input_path, success, error = future.result(timeout=timeout)
                    if success:
                        success_count += 1
                    else:
                        failure_count += 1
                        console.print(
                            f"\nError processing {input_path}: {error}",
                            style="bold red",
                        )
                except TimeoutError:
                    failure_count += 1
                    input_path = future_to_path[future]
                    console.print(
                        f"\nTimeout processing {input_path}: Operation took longer than {timeout} seconds",
                        style="bold red",
                    )
                    future.cancel()
                except Exception as e:
                    failure_count += 1
                    input_path = future_to_path[future]
                    console.print(
                        f"\nUnexpected error processing {input_path}: {str(e)}",
                        style="bold red",
                    )

                progress.update(task, advance=1)

    return success_count, failure_count, skipped_count


def main():
    parser = argparse.ArgumentParser(
        description="Invert colors in PowerPoint presentations"
    )
    parser.add_argument("input", type=Path, help="Input file or directory")
    parser.add_argument(
        "-o", "--output", type=Path, help="Output file or directory (optional)"
    )
    parser.add_argument(
        "-w",
        "--workers",
        type=int,
        help="Number of worker threads (default: CPU count * 4)",
    )
    parser.add_argument(
        "-r",
        "--recursive",
        action="store_true",
        help="Process subdirectories recursively",
    )
    # TODO: This argument does not work correctly
    parser.add_argument(
        "-t",
        "--timeout",
        type=int,
        default=300,
        help="Timeout in seconds for processing each file (default: 300)",
    )
    parser.add_argument(
        "-f",
        "--force",
        action="store_true",
        help="Force processing of all files, even if they exist in output",
    )
    args = parser.parse_args()

    # Validate input path exists
    if not args.input.exists():
        console.print(
            f"Error: Input path does not exist: {args.input}", style="bold red"
        )
        sys.exit(1)

    # If no output specified, create default output directory
    if not args.output:
        args.output = get_default_output_dir(args.input)

    # Process single file
    if args.input.is_file():
        if not args.input.suffix.lower() == ".pptx":
            console.print(
                f"Error: Input file must be a .pptx file: {args.input}",
                style="bold red",
            )
            sys.exit(1)

        # Check if we need to process this file
        if not args.force and args.output.exists():
            existing_files = {args.output: calculate_file_hash(args.output)}
            if not should_process_file(args.input, args.output, existing_files):
                console.print(
                    f"File already processed: {args.output}", style="bold green"
                )
                sys.exit(0)

        console.print(f"Processing: {args.input}", style="bold blue")
        _, success, error = process_file_with_timeout(
            args.input, args.output, args.timeout
        )
        if success:
            console.print(
                f"Successfully processed: {args.input} → {args.output}",
                style="bold green",
            )
        else:
            console.print(f"Error: {error}", style="bold red")
        sys.exit(0 if success else 1)

    # Process directory
    elif args.input.is_dir():
        success_count, failure_count, skipped_count = process_directory(
            args.input,
            args.output,
            recursive=args.recursive,
            max_workers=args.workers,
            timeout=args.timeout,
        )
        total = success_count + failure_count + skipped_count

        if not interrupted:
            console.print("\nProcessing complete:", style="bold blue")
            console.print(
                f"Successfully processed: {success_count}/{total} files",
                style="bold green",
            )
            if failure_count > 0:
                console.print(
                    f"Failed to process: {failure_count}/{total} files",
                    style="bold red",
                )
            if skipped_count > 0:
                console.print(
                    f"Skipped (already processed): {skipped_count}/{total} files",
                    style="bold yellow",
                )
            console.print(f"Output directory: {args.output}", style="bold blue")

        sys.exit(1 if failure_count > 0 or interrupted else 0)


if __name__ == "__main__":
    main()
